from pptx import Presentation as ppt
from util.data_cleaning import separate_glued_words, remove_multiple_whitespaces, merge_string

class Presentation():

  presentation = None
  string = None
  paragraphs = []
  raw_text = None

  def __init__(self, path):
    self.presentation = ppt(path)
    self.raw_text = self.raw_full_text()
    self.string = self.build_string(self.raw_text)
    self.paragraphs = self.get_paragraphs(self.raw_text)

  def raw_full_text(self):
    raw_text = []
    for slide in self.presentation.slides:
      for shape in slide.shapes:
          if not shape.has_text_frame:
              continue
          for paragraph in shape.text_frame.paragraphs:
              for child in paragraph.runs: ## Text run object. Corresponds to a child element in a paragraph.
                  raw_text.append(child.text)
                
    return raw_text

  def trim_text(self, raw_text):
    raw_text = separate_glued_words(raw_text)
    return remove_multiple_whitespaces(raw_text)

  def build_string(self, raw_text):
    return self.trim_text(
        merge_string(raw_text)
    )

  def get_paragraphs(self, raw_text):
    paragraphs = []
    for paragraph in raw_text:
      paragraphs.append(remove_multiple_whitespaces(separate_glued_words(paragraph)))

    return paragraphs


